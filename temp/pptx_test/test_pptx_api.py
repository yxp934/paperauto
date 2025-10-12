"""
测试 python-pptx API 修改PPT内容
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

def create_template_pptx(output_path):
    """创建一个基础模板PPT"""
    print("📝 正在创建基础模板...")

    prs = Presentation()

    # 设置幻灯片尺寸 (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # === 第1页：标题页 ===
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "研究主题标题"
    subtitle.text = "副标题：作者姓名\n日期：2025-01-01"

    # === 第2页：内容页（标题+要点） ===
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide2.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = '研究背景'

    tf = body_shape.text_frame
    tf.text = '主要研究问题'

    p = tf.add_paragraph()
    p.text = '研究目标1'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '研究目标2'
    p.level = 1

    # === 第3页：空白页（我们会手动添加元素） ===
    blank_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide3 = prs.slides.add_slide(blank_slide_layout)

    # 添加标题文本框
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    title_box = slide3.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "自定义布局页面"

    # 设置标题样式
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # 添加内容文本框
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)

    content_box = slide3.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    p = content_frame.paragraphs[0]
    p.text = "这是一个完全自定义的幻灯片"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)

    # 添加形状作为背景装饰
    shape_left = Inches(8)
    shape_top = Inches(0.2)
    shape_width = Inches(1.5)
    shape_height = Inches(1.5)

    shape = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL,
        shape_left, shape_top, shape_width, shape_height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 173, 181)  # 青色圆圈

    # 设置背景颜色
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)  # 深蓝色背景

    # 保存
    prs.save(output_path)
    print(f"✅ 模板创建成功：{output_path}")
    return output_path


def modify_pptx_content(template_path, output_path):
    """修改PPT内容"""
    print(f"\n🔧 正在加载模板：{template_path}")

    prs = Presentation(template_path)

    print(f"📊 模板信息：")
    print(f"   - 幻灯片数量：{len(prs.slides)}")
    print(f"   - 可用布局：{len(prs.slide_layouts)}")
    print(f"   - 幻灯片尺寸：{prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    # === 修改第1页（标题页）===
    print("\n🎨 修改第1页（标题页）...")
    slide1 = prs.slides[0]
    slide1.shapes.title.text = "【已修改】量子计算在机器学习中的应用"
    slide1.placeholders[1].text = "作者：Claude AI\n日期：2025-10-12\n机构：Anthropic Research"

    # === 修改第2页（要点页）===
    print("🎨 修改第2页（要点页）...")
    slide2 = prs.slides[1]
    slide2.shapes.title.text = '【已修改】研究方法论'

    # 清空原有要点，重新添加
    body_shape = slide2.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    # 添加新要点
    new_bullets = [
        ("数据收集与预处理", 0),
        ("量子算法设计", 1),
        ("变分量子特征提取", 2),
        ("经典-量子混合训练", 1),
        ("实验验证与结果分析", 0),
        ("性能对比实验", 1),
        ("误差分析", 1),
    ]

    for i, (text, level) in enumerate(new_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - level * 2)

    # === 修改第3页（自定义页）===
    print("🎨 修改第3页（自定义布局）...")
    slide3 = prs.slides[2]

    # 修改标题
    for shape in slide3.shapes:
        if shape.has_text_frame and "自定义布局页面" in shape.text:
            shape.text = "【已修改】核心技术架构"

    # === 添加新的第4页（带图片）===
    print("🎨 添加第4页（图片+文字）...")
    blank_layout = prs.slide_layouts[6]
    slide4 = prs.slides.add_slide(blank_layout)

    # 设置背景
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)  # 浅灰色背景

    # 添加标题
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "实验结果展示"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 30, 30)
    p.alignment = PP_ALIGN.CENTER

    # 添加描述文本
    desc_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True

    p = desc_frame.paragraphs[0]
    p.text = "关键发现：\n\n"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(60, 60, 60)

    p = desc_frame.add_paragraph()
    p.text = "• 量子算法相比经典算法提升了37%的准确率\n\n"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)

    p = desc_frame.add_paragraph()
    p.text = "• 训练时间缩短了52%\n\n"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)

    p = desc_frame.add_paragraph()
    p.text = "• 能耗降低了68%"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)

    # 添加装饰形状（模拟图表区域）
    chart_area = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.5), Inches(4.5), Inches(3.5)
    )
    fill = chart_area.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(37, 99, 235)  # 蓝色

    # 在"图表"中添加文字
    chart_text = chart_area.text_frame
    chart_text.text = "📊\n\n图表占位符\n（这里可以插入实际图片）"
    for paragraph in chart_text.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # === 添加第5页（结论页）===
    print("🎨 添加第5页（结论）...")
    conclusion_layout = prs.slide_layouts[1]  # Title and Content
    slide5 = prs.slides.add_slide(conclusion_layout)
    slide5.shapes.title.text = "研究结论与展望"

    body = slide5.placeholders[1]
    tf = body.text_frame
    tf.clear()

    conclusions = [
        "量子计算在特定机器学习任务中展现出显著优势",
        "混合算法架构是当前最实用的解决方案",
        "未来研究方向：",
        "扩展到更大规模数据集",
        "探索更多量子算法变体",
        "优化量子-经典接口",
    ]

    for i, text in enumerate(conclusions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0 if i < 3 else 1
        p.font.size = Pt(22 if p.level == 0 else 18)

    # 保存修改后的PPT
    prs.save(output_path)
    print(f"\n✅ PPT修改成功：{output_path}")
    print(f"📊 最终幻灯片数量：{len(prs.slides)}")

    return output_path


def inspect_pptx_structure(pptx_path):
    """检查PPT结构"""
    print(f"\n🔍 检查PPT结构：{pptx_path}")
    print("=" * 60)

    prs = Presentation(pptx_path)

    print(f"\n📐 幻灯片尺寸：{prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"📄 总页数：{len(prs.slides)}")
    print(f"\n🎨 可用布局：")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"   [{i}] {layout.name}")

    print(f"\n📋 各页详情：")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- 第 {i+1} 页 ---")
        print(f"  布局：{slide.slide_layout.name}")
        print(f"  形状数量：{len(slide.shapes)}")

        # 列出所有形状
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_preview = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                print(f"    [{j}] 文本框: {text_preview}")
            elif hasattr(shape, 'image'):
                print(f"    [{j}] 图片")
            else:
                print(f"    [{j}] {shape.shape_type}")


if __name__ == "__main__":
    print("=" * 60)
    print("🧪 python-pptx API 测试程序")
    print("=" * 60)

    # 设置路径
    test_dir = "/Users/yxp/Documents/ghpaperauto/temp/pptx_test"
    template_path = os.path.join(test_dir, "template.pptx")
    modified_path = os.path.join(test_dir, "modified.pptx")

    # 步骤1：创建基础模板
    create_template_pptx(template_path)

    # 步骤2：检查模板结构
    inspect_pptx_structure(template_path)

    # 步骤3：修改内容
    modify_pptx_content(template_path, modified_path)

    # 步骤4：检查修改后的结构
    inspect_pptx_structure(modified_path)

    print("\n" + "=" * 60)
    print("✅ 测试完成！")
    print(f"📁 原始模板：{template_path}")
    print(f"📁 修改后的：{modified_path}")
    print("=" * 60)
    print("\n💡 结论：")
    print("   ✓ python-pptx 可以创建PPT模板")
    print("   ✓ python-pptx 可以修改现有PPT内容")
    print("   ✓ 可以访问和修改标题、要点、文本框")
    print("   ✓ 可以添加新页面、形状、样式")
    print("   ✓ 生成的PPT完全可编辑（非图片）")
